import logging
import io
import zipfile
from typing import Tuple, Optional
import PyPDF2
import openpyxl
from docx import Document
from pptx import Presentation
from utils.logger import setup_logger
from utils.validators import validate_file_upload, get_file_category

logger = setup_logger(__name__)

class FileProcessor:
    def __init__(self):
        self.max_text_length = 10000  # Limit extracted text length
    
    def process_file(self, file_content: bytes, filename: str) -> Tuple[bool, str, Optional[str]]:
        """
        Process uploaded file and extract text content
        Returns: (success, error_message, extracted_text)
        """
        try:
            # Validate file
            is_valid, error_code = validate_file_upload(file_content, filename)
            if not is_valid:
                return False, error_code, None
            
            # Determine file category and process accordingly
            file_category = get_file_category(filename)
            
            if file_category == 'documents':
                return self._process_document(file_content, filename)
            elif file_category == 'spreadsheets':
                return self._process_spreadsheet(file_content, filename)
            elif file_category == 'presentations':
                return self._process_presentation(file_content, filename)
            elif file_category == 'code':
                return self._process_code_file(file_content, filename)
            else:
                return False, 'unsupported_file', None
                
        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")
            return False, 'processing_error', None
    
    def _process_document(self, file_content: bytes, filename: str) -> Tuple[bool, str, Optional[str]]:
        """Process document files (PDF, DOCX, TXT, MD)"""
        try:
            file_ext = filename.lower().split('.')[-1]
            
            if file_ext == 'pdf':
                return self._extract_pdf_text(file_content)
            elif file_ext in ['docx', 'doc']:
                return self._extract_docx_text(file_content)
            elif file_ext in ['txt', 'md']:
                return self._extract_text_file(file_content)
            else:
                return False, 'unsupported_file', None
                
        except Exception as e:
            logger.error(f"Error processing document {filename}: {e}")
            return False, 'processing_error', None
    
    def _process_spreadsheet(self, file_content: bytes, filename: str) -> Tuple[bool, str, Optional[str]]:
        """Process spreadsheet files (XLSX, XLS, CSV)"""
        try:
            file_ext = filename.lower().split('.')[-1]
            
            if file_ext in ['xlsx', 'xls']:
                return self._extract_xlsx_text(file_content)
            elif file_ext == 'csv':
                return self._extract_csv_text(file_content)
            else:
                return False, 'unsupported_file', None
                
        except Exception as e:
            logger.error(f"Error processing spreadsheet {filename}: {e}")
            return False, 'processing_error', None
    
    def _process_presentation(self, file_content: bytes, filename: str) -> Tuple[bool, str, Optional[str]]:
        """Process presentation files (PPTX, PPT)"""
        try:
            file_ext = filename.lower().split('.')[-1]
            
            if file_ext in ['pptx', 'ppt']:
                return self._extract_pptx_text(file_content)
            else:
                return False, 'unsupported_file', None
                
        except Exception as e:
            logger.error(f"Error processing presentation {filename}: {e}")
            return False, 'processing_error', None
    
    def _process_code_file(self, file_content: bytes, filename: str) -> Tuple[bool, str, Optional[str]]:
        """Process code files"""
        try:
            # Try to decode as UTF-8 text
            text_content = file_content.decode('utf-8')
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            return True, '', text_content
            
        except UnicodeDecodeError:
            logger.error(f"Unable to decode code file {filename} as UTF-8")
            return False, 'processing_error', None
    
    def _extract_pdf_text(self, file_content: bytes) -> Tuple[bool, str, Optional[str]]:
        """Extract text from PDF file"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = ""
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return False, 'processing_error', None
    
    def _extract_docx_text(self, file_content: bytes) -> Tuple[bool, str, Optional[str]]:
        """Extract text from DOCX file"""
        try:
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)
            
            text_content = ""
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return False, 'processing_error', None
    
    def _extract_text_file(self, file_content: bytes) -> Tuple[bool, str, Optional[str]]:
        """Extract text from plain text files"""
        try:
            text_content = file_content.decode('utf-8')
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            return True, '', text_content
            
        except UnicodeDecodeError:
            logger.error("Unable to decode text file as UTF-8")
            return False, 'processing_error', None
    
    def _extract_xlsx_text(self, file_content: bytes) -> Tuple[bool, str, Optional[str]]:
        """Extract text from XLSX file"""
        try:
            xlsx_file = io.BytesIO(file_content)
            workbook = openpyxl.load_workbook(xlsx_file, data_only=True)
            
            text_content = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    if any(row_data):  # Skip empty rows
                        text_content += "\t".join(row_data) + "\n"
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {e}")
            return False, 'processing_error', None
    
    def _extract_csv_text(self, file_content: bytes) -> Tuple[bool, str, Optional[str]]:
        """Extract text from CSV file"""
        try:
            text_content = file_content.decode('utf-8')
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            return True, '', text_content
            
        except UnicodeDecodeError:
            logger.error("Unable to decode CSV file as UTF-8")
            return False, 'processing_error', None
    
    def _extract_pptx_text(self, file_content: bytes) -> Tuple[bool, str, Optional[str]]:
        """Extract text from PPTX file"""
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)
            
            text_content = ""
            for i, slide in enumerate(presentation.slides, 1):
                text_content += f"\n--- Slide {i} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return False, 'processing_error', None
