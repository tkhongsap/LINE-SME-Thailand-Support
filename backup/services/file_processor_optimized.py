import logging
import io
import asyncio
import hashlib
import mimetypes
from typing import Tuple, Optional, Dict, Any, AsyncGenerator
from dataclasses import dataclass
from pathlib import Path
import tempfile
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
import time

# Core libraries
import PyPDF2
import openpyxl
from docx import Document
from pptx import Presentation

# Thai language support
import unicodedata
from collections import Counter

from utils.logger import setup_logger
from utils.validators import validate_file_upload, get_file_category
from config import Config

logger = setup_logger(__name__)

@dataclass
class ProcessingResult:
    """Enhanced processing result with metadata"""
    success: bool
    error_code: Optional[str]
    content: Optional[str]
    metadata: Dict[str, Any]
    processing_time: float
    memory_used: int
    content_hash: Optional[str]

@dataclass
class ProcessingOptions:
    """Configuration options for file processing"""
    max_text_length: int = 50000  # Increased from 10k
    enable_streaming: bool = True
    chunk_size: int = 8192
    preserve_formatting: bool = True
    thai_language_optimized: bool = True
    extract_metadata: bool = True
    enable_compression: bool = True

class OptimizedFileProcessor:
    """
    High-performance file processor optimized for Thai SME documents
    Features:
    - Streaming processing for large files
    - Memory-efficient operations
    - Thai language preservation
    - Async processing capabilities
    - Content caching with compression
    - Robust error handling and corruption detection
    """
    
    def __init__(self):
        # Processing configuration
        self.options = ProcessingOptions()
        
        # Thread pool for CPU-intensive operations
        self.executor = ThreadPoolExecutor(max_workers=4, thread_name_prefix="FileProcessor")
        
        # Processing cache with LRU eviction
        self._cache: Dict[str, ProcessingResult] = {}
        self._cache_access_times: Dict[str, float] = {}
        self._max_cache_size = 100
        
        # Performance metrics
        self.metrics = {
            'files_processed': 0,
            'cache_hits': 0,
            'cache_misses': 0,
            'total_processing_time': 0.0,
            'memory_peak': 0
        }
        
        # File type handlers registry
        self._handlers = {
            'pdf': self._process_pdf_optimized,
            'docx': self._process_docx_optimized,
            'doc': self._process_docx_optimized,
            'xlsx': self._process_xlsx_optimized,
            'xls': self._process_xlsx_optimized,
            'pptx': self._process_pptx_optimized,
            'ppt': self._process_pptx_optimized,
            'txt': self._process_text_optimized,
            'md': self._process_text_optimized,
            'csv': self._process_csv_optimized
        }
    
    async def process_file_async(self, file_content: bytes, filename: str, 
                                options: Optional[ProcessingOptions] = None) -> ProcessingResult:
        """
        Asynchronously process file with full optimization features
        """
        start_time = time.time()
        processing_options = options or self.options
        
        try:
            # Generate content hash for caching
            content_hash = hashlib.sha256(file_content).hexdigest()
            
            # Check cache first
            cached_result = self._get_cached_result(content_hash)
            if cached_result:
                self.metrics['cache_hits'] += 1
                logger.info(f"Cache hit for file {filename} (hash: {content_hash[:8]})")
                return cached_result
            
            self.metrics['cache_misses'] += 1
            
            # Validate file
            is_valid, error_code = validate_file_upload(file_content, filename)
            if not is_valid:
                return ProcessingResult(
                    success=False,
                    error_code=error_code,
                    content=None,
                    metadata={},
                    processing_time=time.time() - start_time,
                    memory_used=len(file_content),
                    content_hash=content_hash
                )
            
            # Detect file corruption
            if self._is_file_corrupted(file_content, filename):
                return ProcessingResult(
                    success=False,
                    error_code='file_corrupted',
                    content=None,
                    metadata={},
                    processing_time=time.time() - start_time,
                    memory_used=len(file_content),
                    content_hash=content_hash
                )
            
            # Process file in thread pool for CPU-intensive operations
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(
                self.executor,
                self._process_file_sync,
                file_content,
                filename,
                processing_options,
                content_hash,
                start_time
            )
            
            # Cache successful results
            if result.success:
                self._cache_result(content_hash, result)
            
            # Update metrics
            self.metrics['files_processed'] += 1
            self.metrics['total_processing_time'] += result.processing_time
            
            return result
            
        except Exception as e:
            logger.error(f"Async processing failed for {filename}: {e}")
            return ProcessingResult(
                success=False,
                error_code='processing_error',
                content=None,
                metadata={'error': str(e)},
                processing_time=time.time() - start_time,
                memory_used=len(file_content),
                content_hash=content_hash if 'content_hash' in locals() else None
            )
    
    def process_file(self, file_content: bytes, filename: str) -> Tuple[bool, str, Optional[str]]:
        """
        Synchronous wrapper for backward compatibility
        """
        # Run async processing in new event loop
        try:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            result = loop.run_until_complete(
                self.process_file_async(file_content, filename)
            )
            loop.close()
            
            return result.success, result.error_code or '', result.content
            
        except Exception as e:
            logger.error(f"Sync processing wrapper failed: {e}")
            return False, 'processing_error', None
    
    def _process_file_sync(self, file_content: bytes, filename: str, 
                          options: ProcessingOptions, content_hash: str, 
                          start_time: float) -> ProcessingResult:
        """
        Synchronous file processing with streaming and optimization
        """
        try:
            # Extract file extension
            file_ext = filename.lower().split('.')[-1]
            
            # Get appropriate handler
            handler = self._handlers.get(file_ext)
            if not handler:
                return ProcessingResult(
                    success=False,
                    error_code='unsupported_file',
                    content=None,
                    metadata={},
                    processing_time=time.time() - start_time,
                    memory_used=len(file_content),
                    content_hash=content_hash
                )
            
            # Process with specific handler
            extracted_text, metadata = handler(file_content, filename, options)
            
            # Optimize text content
            if extracted_text:
                extracted_text = self._optimize_text_content(extracted_text, options)
            
            processing_time = time.time() - start_time
            
            result = ProcessingResult(
                success=True,
                error_code=None,
                content=extracted_text,
                metadata=metadata,
                processing_time=processing_time,
                memory_used=len(file_content),
                content_hash=content_hash
            )
            
            logger.info(f"Successfully processed {filename} in {processing_time:.2f}s")
            return result
            
        except Exception as e:
            logger.error(f"File processing failed for {filename}: {e}")
            return ProcessingResult(
                success=False,
                error_code='processing_error',
                content=None,
                metadata={'error': str(e)},
                processing_time=time.time() - start_time,
                memory_used=len(file_content),
                content_hash=content_hash
            )
    
    def _process_pdf_optimized(self, file_content: bytes, filename: str, 
                              options: ProcessingOptions) -> Tuple[str, Dict]:
        """
        Optimized PDF processing with streaming and Thai language support
        """
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text_content = ""
        metadata = {
            'pages': len(pdf_reader.pages),
            'encrypted': pdf_reader.is_encrypted,
            'file_type': 'pdf'
        }
        
        # Process pages with memory-efficient streaming
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                page_text = page.extract_text()
                
                # Thai language optimization - preserve Thai character spacing
                if options.thai_language_optimized:
                    page_text = self._optimize_thai_text(page_text)
                
                text_content += f"\n--- หน้า {page_num + 1} ---\n{page_text}\n"
                
                # Check length limit during processing
                if len(text_content) > options.max_text_length:
                    text_content = text_content[:options.max_text_length] + "\n... (เนื้อหาถูกตัดทอนเนื่องจากมีขนาดใหญ่)"
                    metadata['truncated'] = True
                    break
                    
            except Exception as e:
                logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")
                continue
        
        return text_content.strip(), metadata
    
    def _process_docx_optimized(self, file_content: bytes, filename: str, 
                               options: ProcessingOptions) -> Tuple[str, Dict]:
        """
        Optimized DOCX processing with structure preservation
        """
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)
        
        text_content = ""
        metadata = {
            'paragraphs': len(doc.paragraphs),
            'file_type': 'docx'
        }
        
        # Process paragraphs with formatting preservation
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if para_text:
                # Preserve paragraph structure for Thai documents
                if options.preserve_formatting:
                    text_content += para_text + "\n\n"
                else:
                    text_content += para_text + "\n"
                
                # Check length limit
                if len(text_content) > options.max_text_length:
                    text_content = text_content[:options.max_text_length] + "\n... (เนื้อหาถูกตัดทอนเนื่องจากมีขนาดใหญ่)"
                    metadata['truncated'] = True
                    break
        
        # Process tables if present
        if doc.tables:
            metadata['tables'] = len(doc.tables)
            table_text = self._extract_table_content(doc.tables, options)
            if table_text and len(text_content) + len(table_text) <= options.max_text_length:
                text_content += "\n--- ตาราง ---\n" + table_text
        
        return text_content.strip(), metadata
    
    def _process_xlsx_optimized(self, file_content: bytes, filename: str, 
                               options: ProcessingOptions) -> Tuple[str, Dict]:
        """
        Optimized XLSX processing with Thai number formatting
        """
        xlsx_file = io.BytesIO(file_content)
        workbook = openpyxl.load_workbook(xlsx_file, data_only=True, read_only=True)
        
        text_content = ""
        metadata = {
            'sheets': len(workbook.sheetnames),
            'file_type': 'xlsx'
        }
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content += f"\n--- แผ่นงาน: {sheet_name} ---\n"
            
            # Process rows efficiently
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if row_count > 1000:  # Limit rows to prevent memory issues
                    text_content += "... (แสดงเพียงส่วนหนึ่งของข้อมูล)\n"
                    metadata['truncated'] = True
                    break
                
                row_data = []
                for cell in row:
                    if cell is not None:
                        # Format Thai numbers and currency properly
                        if isinstance(cell, (int, float)):
                            row_data.append(f"{cell:,}")
                        else:
                            row_data.append(str(cell))
                    else:
                        row_data.append('')
                
                if any(row_data):  # Skip empty rows
                    text_content += "\t".join(row_data) + "\n"
                    row_count += 1
                
                # Check length limit
                if len(text_content) > options.max_text_length:
                    text_content = text_content[:options.max_text_length] + "\n... (เนื้อหาถูกตัดทอนเนื่องจากมีขนาดใหญ่)"
                    metadata['truncated'] = True
                    break
        
        workbook.close()
        return text_content.strip(), metadata
    
    def _process_pptx_optimized(self, file_content: bytes, filename: str, 
                               options: ProcessingOptions) -> Tuple[str, Dict]:
        """
        Optimized PPTX processing with slide structure preservation
        """
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        text_content = ""
        metadata = {
            'slides': len(presentation.slides),
            'file_type': 'pptx'
        }
        
        for i, slide in enumerate(presentation.slides, 1):
            text_content += f"\n--- สลไลด์ {i} ---\n"
            
            # Extract text from all shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_content += "\n".join(slide_text) + "\n"
            
            # Check length limit
            if len(text_content) > options.max_text_length:
                text_content = text_content[:options.max_text_length] + "\n... (เนื้อหาถูกตัดทอนเนื่องจากมีขนาดใหญ่)"
                metadata['truncated'] = True
                break
        
        return text_content.strip(), metadata
    
    def _process_text_optimized(self, file_content: bytes, filename: str, 
                               options: ProcessingOptions) -> Tuple[str, Dict]:
        """
        Optimized text file processing with encoding detection
        """
        metadata = {'file_type': 'text'}
        
        # Try multiple encodings for Thai content
        encodings = ['utf-8', 'utf-8-sig', 'tis-620', 'cp874']
        
        for encoding in encodings:
            try:
                text_content = file_content.decode(encoding)
                metadata['encoding'] = encoding
                break
            except UnicodeDecodeError:
                continue
        else:
            # Fallback to latin-1 which never fails
            text_content = file_content.decode('latin-1', errors='replace')
            metadata['encoding'] = 'latin-1'
            metadata['encoding_issues'] = True
        
        # Apply length limit
        if len(text_content) > options.max_text_length:
            text_content = text_content[:options.max_text_length] + "\n... (เนื้อหาถูกตัดทอนเนื่องจากมีขนาดใหญ่)"
            metadata['truncated'] = True
        
        return text_content, metadata
    
    def _process_csv_optimized(self, file_content: bytes, filename: str, 
                              options: ProcessingOptions) -> Tuple[str, Dict]:
        """
        Optimized CSV processing with Thai encoding support
        """
        import csv
        
        # Try multiple encodings for Thai CSV files
        encodings = ['utf-8', 'utf-8-sig', 'tis-620', 'cp874']
        
        for encoding in encodings:
            try:
                text_content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text_content = file_content.decode('latin-1', errors='replace')
        
        metadata = {
            'file_type': 'csv',
            'encoding': encoding if 'encoding' in locals() else 'latin-1'
        }
        
        # Process CSV with streaming for large files
        csv_reader = csv.reader(io.StringIO(text_content))
        processed_content = ""
        row_count = 0
        
        for row in csv_reader:
            if row_count > 1000:  # Limit rows
                processed_content += "... (แสดงเพียงส่วนหนึ่งของข้อมูล)\n"
                metadata['truncated'] = True
                break
            
            if any(row):  # Skip empty rows
                processed_content += "\t".join(row) + "\n"
                row_count += 1
            
            # Check length limit
            if len(processed_content) > options.max_text_length:
                processed_content = processed_content[:options.max_text_length] + "\n... (เนื้อหาถูกตัดทอนเนื่องจากมีขนาดใหญ่)"
                metadata['truncated'] = True
                break
        
        metadata['rows_processed'] = row_count
        return processed_content.strip(), metadata
    
    def _optimize_thai_text(self, text: str) -> str:
        """
        Optimize Thai text by preserving proper character spacing and formatting
        """
        if not text:
            return text
        
        # Normalize Unicode for consistent Thai character representation
        text = unicodedata.normalize('NFC', text)
        
        # Remove excessive whitespace while preserving Thai text structure
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # Preserve Thai text spacing
            cleaned_line = ' '.join(line.split())
            if cleaned_line:
                cleaned_lines.append(cleaned_line)
        
        return '\n'.join(cleaned_lines)
    
    def _optimize_text_content(self, text: str, options: ProcessingOptions) -> str:
        """
        General text content optimization
        """
        if not text:
            return text
        
        if options.thai_language_optimized:
            text = self._optimize_thai_text(text)
        
        # Remove excessive blank lines
        while '\n\n\n' in text:
            text = text.replace('\n\n\n', '\n\n')
        
        return text.strip()
    
    def _extract_table_content(self, tables, options: ProcessingOptions) -> str:
        """
        Extract content from document tables
        """
        table_content = ""
        
        for table_idx, table in enumerate(tables):
            table_content += f"\nตาราง {table_idx + 1}:\n"
            
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_data.append(cell_text)
                
                if row_data:
                    table_content += "\t".join(row_data) + "\n"
        
        return table_content
    
    def _is_file_corrupted(self, file_content: bytes, filename: str) -> bool:
        """
        Detect if file is corrupted based on file type signatures
        """
        if len(file_content) < 4:
            return True
        
        file_ext = filename.lower().split('.')[-1]
        
        # Check file signatures (magic numbers)
        signatures = {
            'pdf': [b'%PDF'],
            'docx': [b'PK\x03\x04'],  # ZIP-based formats
            'xlsx': [b'PK\x03\x04'],
            'pptx': [b'PK\x03\x04'],
        }
        
        if file_ext in signatures:
            for signature in signatures[file_ext]:
                if file_content.startswith(signature):
                    return False
            return True  # No valid signature found
        
        return False  # Assume other formats are OK
    
    def _get_cached_result(self, content_hash: str) -> Optional[ProcessingResult]:
        """
        Get cached processing result
        """
        if content_hash in self._cache:
            # Update access time for LRU
            self._cache_access_times[content_hash] = time.time()
            return self._cache[content_hash]
        return None
    
    def _cache_result(self, content_hash: str, result: ProcessingResult) -> None:
        """
        Cache processing result with LRU eviction
        """
        # Evict oldest entries if cache is full
        if len(self._cache) >= self._max_cache_size:
            # Find oldest entry
            oldest_hash = min(self._cache_access_times.keys(), 
                            key=lambda k: self._cache_access_times[k])
            del self._cache[oldest_hash]
            del self._cache_access_times[oldest_hash]
        
        self._cache[content_hash] = result
        self._cache_access_times[content_hash] = time.time()
    
    def get_metrics(self) -> Dict[str, Any]:
        """
        Get processing performance metrics
        """
        cache_hit_rate = (self.metrics['cache_hits'] / 
                         max(1, self.metrics['cache_hits'] + self.metrics['cache_misses']))
        
        avg_processing_time = (self.metrics['total_processing_time'] / 
                              max(1, self.metrics['files_processed']))
        
        return {
            **self.metrics,
            'cache_hit_rate': cache_hit_rate,
            'avg_processing_time': avg_processing_time,
            'cache_size': len(self._cache)
        }
    
    async def process_batch_async(self, files: List[Tuple[bytes, str]]) -> List[ProcessingResult]:
        """
        Process multiple files concurrently
        """
        tasks = []
        for file_content, filename in files:
            task = self.process_file_async(file_content, filename)
            tasks.append(task)
        
        return await asyncio.gather(*tasks, return_exceptions=True)
    
    def cleanup(self):
        """
        Cleanup resources
        """
        self.executor.shutdown(wait=True)
        self._cache.clear()
        self._cache_access_times.clear()

# Create singleton instance
optimized_file_processor = OptimizedFileProcessor()