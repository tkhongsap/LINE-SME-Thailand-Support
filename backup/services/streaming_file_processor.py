"""
Streaming File Processor with Progressive Loading and Real-time Status Updates
Optimized for Thai SME LINE Bot with enhanced user experience
"""

import logging
import io
import asyncio
import threading
import time
from typing import Tuple, Optional, Dict, Any, Iterator, Callable
from dataclasses import dataclass
from enum import Enum
import json
from concurrent.futures import ThreadPoolExecutor, Future

# File processing libraries
import PyPDF2
import openpyxl
from docx import Document
from pptx import Presentation
# Optional pandas import for CSV processing
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    pd = None

import zipfile
try:
    import rarfile
    RARFILE_AVAILABLE = True
except ImportError:
    RARFILE_AVAILABLE = False
    rarfile = None

import mimetypes

# Image processing
from PIL import Image

# Optional OCR support
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    pytesseract = None

import logging
from utils.validators import validate_file_upload, get_file_category
from config import Config

logger = logging.getLogger(__name__)

class ProcessingStatus(Enum):
    PENDING = "pending"
    VALIDATING = "validating"
    EXTRACTING = "extracting"
    ANALYZING = "analyzing"
    COMPLETED = "completed"
    FAILED = "failed"

@dataclass
class ProcessingProgress:
    status: ProcessingStatus
    progress_percent: int
    current_step: str
    estimated_time_remaining: Optional[float] = None
    extracted_text_length: int = 0
    file_info: Optional[Dict] = None
    error_message: Optional[str] = None

class StreamingFileProcessor:
    """Enhanced file processor with streaming capabilities and real-time progress"""
    
    def __init__(self):
        self.max_text_length = Config.MAX_FILE_SIZE // 100  # Reasonable text limit
        self.chunk_size = 8192  # 8KB chunks for streaming
        self.max_concurrent_processes = 3
        self.executor = ThreadPoolExecutor(max_workers=self.max_concurrent_processes)
        
        # Processing metrics
        self.processing_stats = {
            'files_processed': 0,
            'total_processing_time': 0,
            'success_rate': 0,
            'average_file_size': 0,
            'most_common_types': {}
        }
        
        # OCR configuration (if available)
        self.ocr_enabled = self._check_ocr_availability()
        
        logger.info(f"StreamingFileProcessor initialized with OCR: {self.ocr_enabled}")
    
    def _check_ocr_availability(self) -> bool:
        """Check if OCR (Tesseract) is available"""
        if not TESSERACT_AVAILABLE:
            logger.warning("pytesseract not available - install for image text extraction")
            return False
        
        try:
            pytesseract.get_tesseract_version()
            return True
        except Exception:
            logger.warning("OCR not available - install tesseract for image text extraction")
            return False
    
    async def process_file_streaming(
        self, 
        file_content: bytes, 
        filename: str,
        progress_callback: Optional[Callable[[ProcessingProgress], None]] = None,
        user_context: Optional[Dict] = None
    ) -> Tuple[bool, str, Optional[str], Dict[str, Any]]:
        """
        Process file with streaming progress updates
        
        Returns:
            Tuple[success, error_code, extracted_text, metadata]
        """
        start_time = time.time()
        progress = ProcessingProgress(
            status=ProcessingStatus.PENDING,
            progress_percent=0,
            current_step="Starting file processing..."
        )
        
        # Send initial progress
        if progress_callback:
            progress_callback(progress)
        
        try:
            # Step 1: Validation (10%)
            progress.status = ProcessingStatus.VALIDATING
            progress.progress_percent = 10
            progress.current_step = "Validating file..."
            if progress_callback:
                progress_callback(progress)
            
            is_valid, error_code = validate_file_upload(file_content, filename)
            if not is_valid:
                progress.status = ProcessingStatus.FAILED
                progress.error_message = error_code
                if progress_callback:
                    progress_callback(progress)
                return False, error_code, None, {}
            
            # Step 2: File Analysis (20%)
            progress.progress_percent = 20
            progress.current_step = "Analyzing file structure..."
            if progress_callback:
                progress_callback(progress)
            
            file_info = await self._analyze_file_async(file_content, filename)
            progress.file_info = file_info
            
            # Step 3: Content Extraction (60%)
            progress.status = ProcessingStatus.EXTRACTING
            progress.progress_percent = 30
            progress.current_step = f"Extracting content from {file_info['type']} file..."
            if progress_callback:
                progress_callback(progress)
            
            # Process based on file type with progress updates
            success, error_msg, extracted_text = await self._process_by_type_async(
                file_content, filename, file_info, progress_callback
            )
            
            if not success:
                progress.status = ProcessingStatus.FAILED
                progress.error_message = error_msg
                if progress_callback:
                    progress_callback(progress)
                return False, error_msg, None, file_info
            
            # Step 4: Text Analysis and Optimization (20%)
            progress.status = ProcessingStatus.ANALYZING
            progress.progress_percent = 90
            progress.current_step = "Optimizing extracted content..."
            if progress_callback:
                progress_callback(progress)
            
            # Optimize extracted text
            optimized_text = await self._optimize_text_async(extracted_text, user_context)
            progress.extracted_text_length = len(optimized_text) if optimized_text else 0
            
            # Step 5: Completion (100%)
            progress.status = ProcessingStatus.COMPLETED
            progress.progress_percent = 100
            progress.current_step = "Processing completed successfully!"
            
            processing_time = time.time() - start_time
            progress.estimated_time_remaining = 0
            
            if progress_callback:
                progress_callback(progress)
            
            # Update statistics
            await self._update_processing_stats(processing_time, len(file_content), file_info['type'], True)
            
            # Enhanced metadata
            metadata = {
                **file_info,
                'processing_time_ms': round(processing_time * 1000, 1),
                'text_length': len(optimized_text) if optimized_text else 0,
                'compression_ratio': len(optimized_text) / len(extracted_text) if extracted_text else 1,
                'extraction_method': file_info.get('extraction_method', 'standard')
            }
            
            logger.info(f"Successfully processed {filename} in {processing_time:.2f}s")
            return True, '', optimized_text, metadata
            
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"Error processing file {filename}: {e}")
            
            progress.status = ProcessingStatus.FAILED
            progress.error_message = "processing_error"
            if progress_callback:
                progress_callback(progress)
            
            await self._update_processing_stats(processing_time, len(file_content), 'unknown', False)
            return False, 'processing_error', None, {}
    
    async def _analyze_file_async(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Analyze file structure and metadata asynchronously"""
        def analyze_sync():
            extension = filename.lower().split('.')[-1] if '.' in filename else ''
            category = get_file_category(filename)
            
            file_info = {
                'filename': filename,
                'size_bytes': len(file_content),
                'size_mb': round(len(file_content) / (1024 * 1024), 2),
                'extension': extension,
                'mime_type': mimetypes.guess_type(filename)[0],
                'category': category,
                'encoding': 'binary',
                'type': self._get_file_type_description(extension, category)
            }
            
            # Try to detect encoding for text files
            if file_info['category'] == 'code':
                try:
                    file_content.decode('utf-8')
                    file_info['encoding'] = 'utf-8'
                except UnicodeDecodeError:
                    try:
                        file_content.decode('windows-1252')
                        file_info['encoding'] = 'windows-1252'
                    except UnicodeDecodeError:
                        file_info['encoding'] = 'unknown'
            
            # Detailed analysis by type
            if file_info['extension'] == 'pdf':
                file_info.update(self._analyze_pdf_structure(file_content))
            elif file_info['extension'] in ['xlsx', 'xls']:
                file_info.update(self._analyze_excel_structure(file_content))
            elif file_info['extension'] == 'docx':
                file_info.update(self._analyze_docx_structure(file_content))
            elif file_info['extension'] == 'pptx':
                file_info.update(self._analyze_pptx_structure(file_content))
            
            return file_info
        
        # Run analysis in thread pool to avoid blocking
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(self.executor, analyze_sync)
    
    def _get_file_type_description(self, extension: str, category: str) -> str:
        """Get human-readable file type description"""
        type_descriptions = {
            'pdf': 'PDF Document',
            'docx': 'Word Document', 
            'doc': 'Word Document',
            'xlsx': 'Excel Spreadsheet',
            'xls': 'Excel Spreadsheet',
            'pptx': 'PowerPoint Presentation',
            'ppt': 'PowerPoint Presentation',
            'txt': 'Text File',
            'md': 'Markdown Document',
            'csv': 'CSV Spreadsheet',
            'py': 'Python Code',
            'js': 'JavaScript Code',
            'html': 'HTML Document',
            'css': 'CSS Stylesheet',
            'java': 'Java Code',
            'cpp': 'C++ Code',
            'c': 'C Code',
            'jpg': 'JPEG Image',
            'jpeg': 'JPEG Image',
            'png': 'PNG Image',
            'gif': 'GIF Image',
            'webp': 'WebP Image'
        }
        
        return type_descriptions.get(extension, f'{category.title()} File')
    
    def _analyze_pdf_structure(self, file_content: bytes) -> Dict[str, Any]:
        """Analyze PDF structure"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            return {
                'type': 'PDF Document',
                'pages': len(pdf_reader.pages),
                'has_text': any(page.extract_text().strip() for page in pdf_reader.pages[:3]),
                'metadata': pdf_reader.metadata if hasattr(pdf_reader, 'metadata') else {},
                'extraction_method': 'text' if any(page.extract_text().strip() for page in pdf_reader.pages[:3]) else 'ocr_required'
            }
        except Exception as e:
            logger.warning(f"PDF analysis failed: {e}")
            return {'type': 'PDF Document', 'analysis_error': str(e)}
    
    def _analyze_excel_structure(self, file_content: bytes) -> Dict[str, Any]:
        """Analyze Excel structure"""
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True)
            sheets_info = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheets_info.append({
                    'name': sheet_name,
                    'max_row': sheet.max_row,
                    'max_column': sheet.max_column,
                    'has_data': sheet.max_row > 1 or sheet.max_column > 1
                })
            
            workbook.close()
            return {
                'type': 'Excel Spreadsheet',
                'sheets': len(sheets_info),
                'sheets_info': sheets_info,
                'extraction_method': 'structured_data'
            }
        except Exception as e:
            logger.warning(f"Excel analysis failed: {e}")
            return {'type': 'Excel Spreadsheet', 'analysis_error': str(e)}
    
    def _analyze_docx_structure(self, file_content: bytes) -> Dict[str, Any]:
        """Analyze DOCX structure"""
        try:
            doc = Document(io.BytesIO(file_content))
            
            return {
                'type': 'Word Document',
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'has_images': any(rel.target_ref.endswith(('.png', '.jpg', '.jpeg', '.gif')) 
                                for rel in doc.part.rels.values() if hasattr(rel, 'target_ref')),
                'extraction_method': 'structured_text'
            }
        except Exception as e:
            logger.warning(f"DOCX analysis failed: {e}")
            return {'type': 'Word Document', 'analysis_error': str(e)}
    
    def _analyze_pptx_structure(self, file_content: bytes) -> Dict[str, Any]:
        """Analyze PPTX structure"""
        try:
            ppt = Presentation(io.BytesIO(file_content))
            
            return {
                'type': 'PowerPoint Presentation',
                'slides': len(ppt.slides),
                'has_text': any(hasattr(shape, 'text') for slide in ppt.slides for shape in slide.shapes),
                'extraction_method': 'slide_text'
            }
        except Exception as e:
            logger.warning(f"PPTX analysis failed: {e}")
            return {'type': 'PowerPoint Presentation', 'analysis_error': str(e)}
    
    async def _process_by_type_async(
        self, 
        file_content: bytes, 
        filename: str, 
        file_info: Dict[str, Any],
        progress_callback: Optional[Callable] = None
    ) -> Tuple[bool, str, Optional[str]]:
        """Process file by type with progress updates"""
        
        category = file_info.get('category', 'unknown')
        extension = file_info.get('extension', '')
        
        def update_progress(percent: int, step: str):
            if progress_callback:
                progress = ProcessingProgress(
                    status=ProcessingStatus.EXTRACTING,
                    progress_percent=30 + int(percent * 0.6),  # 30-90% range
                    current_step=step
                )
                progress_callback(progress)
        
        try:
            if category == 'documents':
                return await self._process_document_streaming(file_content, filename, extension, update_progress)
            elif category == 'spreadsheets':
                return await self._process_spreadsheet_streaming(file_content, filename, extension, update_progress)
            elif category == 'presentations':
                return await self._process_presentation_streaming(file_content, filename, extension, update_progress)
            elif category == 'code':
                return await self._process_code_streaming(file_content, filename, update_progress)
            elif category == 'images':
                return await self._process_image_streaming(file_content, filename, update_progress)
            else:
                return False, 'unsupported_file', None
                
        except Exception as e:
            logger.error(f"Error in type-specific processing: {e}")
            return False, 'processing_error', None
    
    async def _process_document_streaming(
        self, 
        file_content: bytes, 
        filename: str, 
        extension: str,
        progress_callback: Callable
    ) -> Tuple[bool, str, Optional[str]]:
        """Process documents with streaming progress"""
        
        def process_sync():
            if extension == 'pdf':
                progress_callback(25, "Reading PDF pages...")
                return self._extract_pdf_streaming(file_content, progress_callback)
            elif extension in ['docx']:
                progress_callback(25, "Processing Word document...")
                return self._extract_docx_streaming(file_content, progress_callback)
            elif extension in ['txt', 'md']:
                progress_callback(25, "Reading text file...")
                return self._extract_text_file_streaming(file_content, progress_callback)
            else:
                return False, 'unsupported_file', None
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(self.executor, process_sync)
    
    def _extract_pdf_streaming(self, file_content: bytes, progress_callback: Callable) -> Tuple[bool, str, Optional[str]]:
        """Extract PDF text with progress updates"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            total_pages = len(pdf_reader.pages)
            
            text_content = ""
            
            for i, page in enumerate(pdf_reader.pages):
                progress_percent = (i / total_pages) * 100
                progress_callback(progress_percent, f"Processing page {i+1} of {total_pages}...")
                
                page_text = page.extract_text()
                text_content += page_text + "\n"
                
                # Early exit if text is getting too long
                if len(text_content) > self.max_text_length:
                    text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
                    break
            
            progress_callback(100, "PDF processing completed")
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return False, 'processing_error', None
    
    def _extract_docx_streaming(self, file_content: bytes, progress_callback: Callable) -> Tuple[bool, str, Optional[str]]:
        """Extract DOCX text with progress updates"""
        try:
            doc = Document(io.BytesIO(file_content))
            total_paragraphs = len(doc.paragraphs)
            
            text_content = ""
            
            for i, paragraph in enumerate(doc.paragraphs):
                if i % 10 == 0:  # Update progress every 10 paragraphs
                    progress_percent = (i / total_paragraphs) * 100
                    progress_callback(progress_percent, f"Processing paragraph {i+1} of {total_paragraphs}...")
                
                text_content += paragraph.text + "\n"
                
                # Early exit if text is getting too long
                if len(text_content) > self.max_text_length:
                    text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
                    break
            
            # Process tables if they exist
            if doc.tables:
                progress_callback(80, "Processing tables...")
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text_content += cell.text + " "
                        text_content += "\n"
                        
                        if len(text_content) > self.max_text_length:
                            text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
                            break
            
            progress_callback(100, "Document processing completed")
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return False, 'processing_error', None
    
    def _extract_text_file_streaming(self, file_content: bytes, progress_callback: Callable) -> Tuple[bool, str, Optional[str]]:
        """Extract text file content with encoding detection"""
        try:
            progress_callback(25, "Detecting file encoding...")
            
            # Try different encodings
            encodings = ['utf-8', 'windows-1252', 'iso-8859-1', 'cp1251']
            text_content = None
            
            for encoding in encodings:
                try:
                    text_content = file_content.decode(encoding)
                    progress_callback(50, f"Successfully decoded with {encoding}")
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                # Fallback to binary reading with replacement
                text_content = file_content.decode('utf-8', errors='replace')
                progress_callback(75, "Using fallback encoding with replacement")
            
            # Limit text length
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            progress_callback(100, "Text extraction completed")
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting text file: {e}")
            return False, 'processing_error', None
    
    async def _process_spreadsheet_streaming(
        self, 
        file_content: bytes, 
        filename: str, 
        extension: str,
        progress_callback: Callable
    ) -> Tuple[bool, str, Optional[str]]:
        """Process spreadsheets with streaming progress"""
        
        def process_sync():
            if extension in ['xlsx', 'xls']:
                progress_callback(25, "Loading spreadsheet...")
                return self._extract_xlsx_streaming(file_content, progress_callback)
            elif extension == 'csv':
                progress_callback(25, "Reading CSV file...")
                return self._extract_csv_streaming(file_content, progress_callback)
            else:
                return False, 'unsupported_file', None
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(self.executor, process_sync)
    
    def _extract_xlsx_streaming(self, file_content: bytes, progress_callback: Callable) -> Tuple[bool, str, Optional[str]]:
        """Extract Excel data with progress updates"""
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(file_content), read_only=True, data_only=True)
            sheet_names = workbook.sheetnames
            total_sheets = len(sheet_names)
            
            text_content = ""
            
            for sheet_idx, sheet_name in enumerate(sheet_names):
                progress_percent = (sheet_idx / total_sheets) * 100
                progress_callback(progress_percent, f"Processing sheet '{sheet_name}' ({sheet_idx+1}/{total_sheets})...")
                
                sheet = workbook[sheet_name]
                text_content += f"\n=== Sheet: {sheet_name} ===\n"
                
                # Process rows in chunks
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    if row_count > 1000:  # Limit rows processed
                        text_content += "\n... (additional rows truncated)\n"
                        break
                    
                    # Filter out None values and convert to strings
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    if any(row_data):  # Only add non-empty rows
                        text_content += "\t".join(row_data) + "\n"
                    
                    row_count += 1
                    
                    # Check length limit
                    if len(text_content) > self.max_text_length:
                        text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
                        workbook.close()
                        progress_callback(100, "Spreadsheet processing completed (truncated)")
                        return True, '', text_content.strip()
            
            workbook.close()
            progress_callback(100, "Spreadsheet processing completed")
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting Excel text: {e}")
            return False, 'processing_error', None
    
    def _extract_csv_streaming(self, file_content: bytes, progress_callback: Callable) -> Tuple[bool, str, Optional[str]]:
        """Extract CSV data with progress updates"""
        try:
            progress_callback(25, "Reading CSV data...")
            
            # Try different encodings
            encodings = ['utf-8', 'windows-1252', 'iso-8859-1']
            df = None
            
            if PANDAS_AVAILABLE:
                for encoding in encodings:
                    try:
                        df = pd.read_csv(io.BytesIO(file_content), encoding=encoding, nrows=1000)  # Limit rows
                        progress_callback(50, f"Successfully parsed CSV with {encoding}")
                        break
                    except UnicodeDecodeError:
                        continue
                    except Exception as e:
                        logger.warning(f"Error reading CSV with {encoding}: {e}")
                
                if df is None:
                    return False, 'processing_error', None
                
                progress_callback(75, "Converting to text format...")
                
                # Convert DataFrame to text
                text_content = f"CSV Data ({len(df)} rows, {len(df.columns)} columns):\n\n"
                text_content += df.to_string(max_rows=500, max_cols=20)  # Limit output
            else:
                # Fallback CSV processing without pandas
                progress_callback(50, "Processing CSV without pandas (basic mode)")
                
                # Try different encodings for basic text reading
                csv_text = None
                for encoding in encodings:
                    try:
                        csv_text = file_content.decode(encoding)
                        break
                    except UnicodeDecodeError:
                        continue
                
                if csv_text is None:
                    return False, 'processing_error', None
                
                progress_callback(75, "Converting to text format...")
                
                # Basic CSV formatting
                lines = csv_text.split('\n')[:1000]  # Limit lines
                text_content = f"CSV Data ({len(lines)} lines):\n\n"
                text_content += '\n'.join(lines)
            
            if len(text_content) > self.max_text_length:
                text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
            
            progress_callback(100, "CSV processing completed")
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting CSV text: {e}")
            return False, 'processing_error', None
    
    async def _process_presentation_streaming(
        self, 
        file_content: bytes, 
        filename: str, 
        extension: str,
        progress_callback: Callable
    ) -> Tuple[bool, str, Optional[str]]:
        """Process presentations with streaming progress"""
        
        def process_sync():
            if extension in ['pptx', 'ppt']:
                progress_callback(25, "Loading presentation...")
                return self._extract_pptx_streaming(file_content, progress_callback)
            else:
                return False, 'unsupported_file', None
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(self.executor, process_sync)
    
    def _extract_pptx_streaming(self, file_content: bytes, progress_callback: Callable) -> Tuple[bool, str, Optional[str]]:
        """Extract PowerPoint text with progress updates"""
        try:
            ppt = Presentation(io.BytesIO(file_content))
            total_slides = len(ppt.slides)
            
            text_content = f"PowerPoint Presentation ({total_slides} slides):\n\n"
            
            for slide_idx, slide in enumerate(ppt.slides):
                progress_percent = (slide_idx / total_slides) * 100
                progress_callback(progress_percent, f"Processing slide {slide_idx+1} of {total_slides}...")
                
                text_content += f"=== Slide {slide_idx + 1} ===\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text_content += shape.text + "\n"
                
                text_content += "\n"
                
                # Check length limit
                if len(text_content) > self.max_text_length:
                    text_content = text_content[:self.max_text_length] + "\n... (content truncated)"
                    break
            
            progress_callback(100, "Presentation processing completed")
            return True, '', text_content.strip()
            
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return False, 'processing_error', None
    
    async def _process_code_streaming(
        self, 
        file_content: bytes, 
        filename: str,
        progress_callback: Callable
    ) -> Tuple[bool, str, Optional[str]]:
        """Process code files with streaming progress"""
        
        def process_sync():
            progress_callback(25, "Detecting code file encoding...")
            
            # Try different encodings
            encodings = ['utf-8', 'windows-1252', 'iso-8859-1']
            text_content = None
            
            for encoding in encodings:
                try:
                    text_content = file_content.decode(encoding)
                    progress_callback(50, f"Successfully decoded code with {encoding}")
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                return False, 'processing_error', None
            
            progress_callback(75, "Processing code content...")
            
            # Add syntax highlighting hints for AI
            extension = filename.lower().split('.')[-1] if '.' in filename else ''
            language_map = {
                'py': 'Python', 'js': 'JavaScript', 'java': 'Java', 'cpp': 'C++',
                'c': 'C', 'php': 'PHP', 'rb': 'Ruby', 'go': 'Go', 'rs': 'Rust',
                'html': 'HTML', 'css': 'CSS', 'sql': 'SQL', 'json': 'JSON'
            }
            
            language = language_map.get(extension, 'Code')
            formatted_content = f"{language} Code ({filename}):\n\n```{extension}\n{text_content}\n```"
            
            # Limit length
            if len(formatted_content) > self.max_text_length:
                formatted_content = formatted_content[:self.max_text_length] + "\n... (content truncated)"
            
            progress_callback(100, "Code processing completed")
            return True, '', formatted_content.strip()
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(self.executor, process_sync)
    
    async def _process_image_streaming(
        self, 
        file_content: bytes, 
        filename: str,
        progress_callback: Callable
    ) -> Tuple[bool, str, Optional[str]]:
        """Process images with OCR if available"""
        
        def process_sync():
            if not self.ocr_enabled:
                progress_callback(100, "Image uploaded (OCR not available)")
                return True, '', f"Image file uploaded: {filename}\n(Text extraction from images requires OCR setup)"
            
            try:
                progress_callback(25, "Loading image...")
                
                # Open image
                image = Image.open(io.BytesIO(file_content))
                
                progress_callback(50, "Performing OCR text extraction...")
                
                # Extract text using OCR
                if TESSERACT_AVAILABLE:
                    extracted_text = pytesseract.image_to_string(image, lang='eng+tha')  # English + Thai
                else:
                    extracted_text = "OCR not available - install pytesseract and tesseract"
                
                progress_callback(75, "Processing extracted text...")
                
                if extracted_text.strip():
                    formatted_content = f"Text extracted from image ({filename}):\n\n{extracted_text}"
                    
                    if len(formatted_content) > self.max_text_length:
                        formatted_content = formatted_content[:self.max_text_length] + "\n... (content truncated)"
                    
                    progress_callback(100, "Image OCR completed")
                    return True, '', formatted_content.strip()
                else:
                    progress_callback(100, "No text found in image")
                    return True, '', f"Image file uploaded: {filename}\n(No readable text found in image)"
                
            except Exception as e:
                logger.error(f"Error processing image: {e}")
                return False, 'processing_error', None
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(self.executor, process_sync)
    
    async def _optimize_text_async(self, text: str, user_context: Optional[Dict] = None) -> str:
        """Optimize extracted text for AI processing"""
        if not text:
            return text
        
        def optimize_sync():
            # Remove excessive whitespace
            optimized = ' '.join(text.split())
            
            # Remove very long lines that might be formatting artifacts
            lines = optimized.split('\n')
            filtered_lines = []
            
            for line in lines:
                if len(line) < 500:  # Skip very long lines (likely formatting issues)
                    filtered_lines.append(line)
                elif len(line) > 500:
                    # Truncate very long lines
                    filtered_lines.append(line[:500] + "...")
            
            optimized = '\n'.join(filtered_lines)
            
            # Add context for Thai SME analysis
            if user_context and user_context.get('language') == 'th':
                optimized = f"[เอกสารสำหรับ SME ไทย]\n\n{optimized}"
            
            return optimized
        
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(self.executor, optimize_sync)
    
    async def _update_processing_stats(self, processing_time: float, file_size: int, file_type: str, success: bool):
        """Update processing statistics asynchronously"""
        self.processing_stats['files_processed'] += 1
        self.processing_stats['total_processing_time'] += processing_time
        
        if success:
            current_success_rate = self.processing_stats['success_rate']
            total_files = self.processing_stats['files_processed']
            successful_files = int(current_success_rate * (total_files - 1) / 100) + 1
            self.processing_stats['success_rate'] = (successful_files / total_files) * 100
        
        # Update average file size
        current_avg = self.processing_stats['average_file_size']
        total_files = self.processing_stats['files_processed']
        self.processing_stats['average_file_size'] = (
            (current_avg * (total_files - 1) + file_size) / total_files
        )
        
        # Update file type statistics
        if file_type not in self.processing_stats['most_common_types']:
            self.processing_stats['most_common_types'][file_type] = 0
        self.processing_stats['most_common_types'][file_type] += 1
    
    def get_processing_metrics(self) -> Dict[str, Any]:
        """Get comprehensive processing metrics"""
        total_files = self.processing_stats['files_processed']
        
        return {
            'performance': {
                'total_files_processed': total_files,
                'success_rate_percent': round(self.processing_stats['success_rate'], 1),
                'average_processing_time_ms': round(
                    (self.processing_stats['total_processing_time'] / max(1, total_files)) * 1000, 1
                ),
                'average_file_size_mb': round(
                    self.processing_stats['average_file_size'] / (1024 * 1024), 2
                ),
                'concurrent_processes': self.max_concurrent_processes,
                'ocr_enabled': self.ocr_enabled
            },
            'file_types': dict(sorted(
                self.processing_stats['most_common_types'].items(),
                key=lambda x: x[1],
                reverse=True
            )),
            'capabilities': {
                'supported_formats': [
                    'PDF', 'DOCX', 'XLSX', 'PPTX', 'TXT', 'MD', 'CSV',
                    'Python', 'JavaScript', 'Java', 'C++', 'HTML', 'CSS'
                ],
                'streaming_processing': True,
                'progress_tracking': True,
                'ocr_support': self.ocr_enabled,
                'max_file_size_mb': Config.MAX_FILE_SIZE / (1024 * 1024)
            }
        }

# Global streaming processor instance
streaming_processor = StreamingFileProcessor()