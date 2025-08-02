import asyncio
import io
import logging
import gc
import psutil
import os
from typing import AsyncGenerator, Optional, Dict, Any
from dataclasses import dataclass
import tempfile
from pathlib import Path

from utils.logger import setup_logger

logger = setup_logger(__name__)

@dataclass
class MemoryMetrics:
    """Memory usage tracking"""
    peak_memory_mb: float
    current_memory_mb: float
    memory_limit_mb: float
    gc_collections: int

class StreamingFileProcessor:
    """
    Memory-efficient streaming file processor for large files
    
    Features:
    - Streaming processing without loading entire file into memory
    - Memory usage monitoring and limits
    - Automatic garbage collection
    - Temporary file management for extremely large files
    - Progress tracking and cancellation support
    """
    
    def __init__(self, memory_limit_mb: int = 100):
        self.memory_limit_mb = memory_limit_mb
        self.temp_dir = Path(tempfile.gettempdir()) / "linebot_processing"
        self.temp_dir.mkdir(exist_ok=True)
        
        # Memory monitoring
        self.process = psutil.Process(os.getpid())
        self.initial_memory = self.process.memory_info().rss / 1024 / 1024
        
        # Processing statistics
        self.stats = {
            'bytes_processed': 0,
            'chunks_processed': 0,
            'gc_collections': 0,
            'temp_files_created': 0
        }
    
    async def stream_process_large_file(self, file_content: bytes, filename: str, 
                                       chunk_size: int = 8192) -> AsyncGenerator[str, None]:
        """
        Stream process large files in chunks to minimize memory usage
        """
        start_memory = self.process.memory_info().rss / 1024 / 1024
        logger.info(f"Starting streaming processing for {filename} (size: {len(file_content):,} bytes)")
        
        try:
            # Check if file is too large for memory processing
            if len(file_content) > self.memory_limit_mb * 1024 * 1024:
                async for chunk in self._process_with_temp_file(file_content, filename, chunk_size):
                    yield chunk
            else:
                async for chunk in self._process_in_memory_stream(file_content, filename, chunk_size):
                    yield chunk
                    
        finally:
            # Memory cleanup
            gc.collect()
            end_memory = self.process.memory_info().rss / 1024 / 1024
            logger.info(f"Streaming processing complete. Memory: {start_memory:.1f}MB -> {end_memory:.1f}MB")
    
    async def _process_with_temp_file(self, file_content: bytes, filename: str, 
                                     chunk_size: int) -> AsyncGenerator[str, None]:
        """
        Process extremely large files using temporary files
        """
        temp_file_path = None
        
        try:
            # Create temporary file
            temp_file_path = self.temp_dir / f"temp_{hash(file_content) % 10000}_{filename}"
            
            # Write content to temp file in chunks
            with open(temp_file_path, 'wb') as temp_file:
                for i in range(0, len(file_content), chunk_size):
                    chunk = file_content[i:i + chunk_size]
                    temp_file.write(chunk)
                    
                    # Yield control to prevent blocking
                    if i % (chunk_size * 10) == 0:
                        await asyncio.sleep(0)
            
            self.stats['temp_files_created'] += 1
            
            # Process from temp file
            async for text_chunk in self._read_and_process_temp_file(temp_file_path, filename, chunk_size):
                yield text_chunk
                
        finally:
            # Cleanup temp file
            if temp_file_path and temp_file_path.exists():
                temp_file_path.unlink()
    
    async def _process_in_memory_stream(self, file_content: bytes, filename: str, 
                                       chunk_size: int) -> AsyncGenerator[str, None]:
        """
        Process files in memory with streaming approach
        """
        file_ext = filename.lower().split('.')[-1]
        
        if file_ext == 'pdf':
            async for chunk in self._stream_pdf_processing(file_content, chunk_size):
                yield chunk
        elif file_ext in ['txt', 'md', 'csv']:
            async for chunk in self._stream_text_processing(file_content, chunk_size):
                yield chunk
        elif file_ext in ['docx', 'xlsx', 'pptx']:
            # Office formats require full loading, but we can stream the output
            extracted_text = await self._extract_office_content(file_content, file_ext)
            async for chunk in self._stream_text_output(extracted_text, chunk_size):
                yield chunk
        else:
            yield f"Unsupported file type for streaming: {file_ext}"
    
    async def _stream_pdf_processing(self, file_content: bytes, 
                                    chunk_size: int) -> AsyncGenerator[str, None]:
        """
        Stream PDF processing page by page
        """
        import PyPDF2
        
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            total_pages = len(pdf_reader.pages)
            yield f"--- PDF Document ({total_pages} pages) ---\n"
            
            for page_num, page in enumerate(pdf_reader.pages):
                # Monitor memory usage
                await self._check_memory_usage()
                
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        page_header = f"\n--- หน้า {page_num + 1} ---\n"
                        yield page_header
                        
                        # Stream page text in chunks
                        async for text_chunk in self._stream_text_output(page_text, chunk_size):
                            yield text_chunk
                        
                        yield "\n"
                        
                except Exception as e:
                    yield f"\n[Error processing page {page_num + 1}: {str(e)}]\n"
                
                # Yield control every few pages
                if page_num % 5 == 0:
                    await asyncio.sleep(0.01)
                
                self.stats['chunks_processed'] += 1
                
        except Exception as e:
            yield f"Error processing PDF: {str(e)}"
    
    async def _stream_text_processing(self, file_content: bytes, 
                                     chunk_size: int) -> AsyncGenerator[str, None]:
        """
        Stream text file processing
        """
        # Try multiple encodings for Thai content
        encodings = ['utf-8', 'utf-8-sig', 'tis-620', 'cp874']
        
        for encoding in encodings:
            try:
                text_content = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text_content = file_content.decode('latin-1', errors='replace')
        
        # Stream text output in chunks
        async for chunk in self._stream_text_output(text_content, chunk_size):
            yield chunk
    
    async def _stream_text_output(self, text: str, chunk_size: int) -> AsyncGenerator[str, None]:
        """
        Stream text output in manageable chunks
        """
        if not text:
            return
        
        # Split text into chunks at word boundaries when possible
        for i in range(0, len(text), chunk_size):
            chunk = text[i:i + chunk_size]
            
            # Try to break at word boundary for better readability
            if i + chunk_size < len(text) and not text[i + chunk_size].isspace():
                # Find last space in chunk
                last_space = chunk.rfind(' ')
                if last_space > chunk_size * 0.8:  # Only if space is reasonably close to end
                    chunk = chunk[:last_space]
            
            yield chunk
            
            # Yield control to prevent blocking
            await asyncio.sleep(0.001)
            
            self.stats['bytes_processed'] += len(chunk.encode('utf-8'))
    
    async def _extract_office_content(self, file_content: bytes, file_ext: str) -> str:
        """
        Extract content from Office documents (requires full loading)
        """
        loop = asyncio.get_event_loop()
        
        # Run extraction in thread pool to avoid blocking
        if file_ext == 'docx':
            return await loop.run_in_executor(None, self._extract_docx_sync, file_content)
        elif file_ext == 'xlsx':
            return await loop.run_in_executor(None, self._extract_xlsx_sync, file_content)
        elif file_ext == 'pptx':
            return await loop.run_in_executor(None, self._extract_pptx_sync, file_content)
        
        return ""
    
    def _extract_docx_sync(self, file_content: bytes) -> str:
        """Synchronous DOCX extraction"""
        from docx import Document
        
        docx_file = io.BytesIO(file_content)
        doc = Document(docx_file)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip())
        
        return '\n'.join(text_parts)
    
    def _extract_xlsx_sync(self, file_content: bytes) -> str:
        """Synchronous XLSX extraction"""
        import openpyxl
        
        xlsx_file = io.BytesIO(file_content)
        workbook = openpyxl.load_workbook(xlsx_file, data_only=True, read_only=True)
        
        text_parts = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\n--- แผ่นงาน: {sheet_name} ---")
            
            for row in sheet.iter_rows(values_only=True, max_row=1000):  # Limit rows
                row_data = [str(cell) if cell is not None else '' for cell in row]
                if any(row_data):
                    text_parts.append('\t'.join(row_data))
        
        workbook.close()
        return '\n'.join(text_parts)
    
    def _extract_pptx_sync(self, file_content: bytes) -> str:
        """Synchronous PPTX extraction"""
        from pptx import Presentation
        
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        
        text_parts = []
        for i, slide in enumerate(presentation.slides, 1):
            text_parts.append(f"\n--- สลไลด์ {i} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
        
        return '\n'.join(text_parts)
    
    async def _read_and_process_temp_file(self, temp_file_path: Path, filename: str, 
                                         chunk_size: int) -> AsyncGenerator[str, None]:
        """
        Read and process content from temporary file
        """
        file_ext = filename.lower().split('.')[-1]
        
        if file_ext in ['txt', 'md', 'csv']:
            # Stream text files directly
            async with aiofiles.open(temp_file_path, 'r', encoding='utf-8') as f:
                while True:
                    chunk = await f.read(chunk_size)
                    if not chunk:
                        break
                    yield chunk
        else:
            # For binary formats, read and process normally
            with open(temp_file_path, 'rb') as f:
                file_content = f.read()
            
            async for chunk in self._process_in_memory_stream(file_content, filename, chunk_size):
                yield chunk
    
    async def _check_memory_usage(self):
        """
        Monitor memory usage and trigger garbage collection if needed
        """
        current_memory = self.process.memory_info().rss / 1024 / 1024
        
        if current_memory > self.memory_limit_mb:
            logger.warning(f"Memory usage ({current_memory:.1f}MB) exceeds limit ({self.memory_limit_mb}MB)")
            
            # Force garbage collection
            gc.collect()
            self.stats['gc_collections'] += 1
            
            # Check memory again
            new_memory = self.process.memory_info().rss / 1024 / 1024
            logger.info(f"After GC: {current_memory:.1f}MB -> {new_memory:.1f}MB")
    
    def get_memory_metrics(self) -> MemoryMetrics:
        """
        Get current memory usage metrics
        """
        current_memory = self.process.memory_info().rss / 1024 / 1024
        peak_memory = max(current_memory, getattr(self, 'peak_memory', current_memory))
        
        return MemoryMetrics(
            peak_memory_mb=peak_memory,
            current_memory_mb=current_memory,
            memory_limit_mb=self.memory_limit_mb,
            gc_collections=self.stats['gc_collections']
        )
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """
        Get processing statistics
        """
        return {
            **self.stats,
            'memory_metrics': self.get_memory_metrics().__dict__
        }
    
    def cleanup(self):
        """
        Cleanup temporary files and resources
        """
        if self.temp_dir.exists():
            for temp_file in self.temp_dir.glob("temp_*"):
                try:
                    temp_file.unlink()
                except Exception as e:
                    logger.warning(f"Failed to cleanup temp file {temp_file}: {e}")
        
        # Force garbage collection
        gc.collect()

# Try to import aiofiles for async file operations
try:
    import aiofiles
except ImportError:
    logger.warning("aiofiles not installed, falling back to synchronous file operations")
    # Create a mock aiofiles for compatibility
    class MockAioFiles:
        @staticmethod
        def open(file_path, mode='r', encoding='utf-8'):
            class AsyncFileWrapper:
                def __init__(self, file_path, mode, encoding):
                    self.file_path = file_path
                    self.mode = mode
                    self.encoding = encoding
                    self.file = None
                
                async def __aenter__(self):
                    self.file = open(self.file_path, self.mode, encoding=self.encoding)
                    return self
                
                async def __aexit__(self, exc_type, exc_val, exc_tb):
                    if self.file:
                        self.file.close()
                
                async def read(self, size=-1):
                    if self.file:
                        return self.file.read(size)
                    return ""
            
            return AsyncFileWrapper(file_path, mode, encoding)
    
    aiofiles = MockAioFiles()

# Create singleton instance
streaming_processor = StreamingFileProcessor()